from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("K:/projects/cadgaurd/CADGuard_AI_Intelligent_Design_Validation_Platform_v2.pptx")
ASSET_DIR = Path("K:/projects/cadgaurd/scripts/ppt_assets")

BG_DARK = RGBColor(5, 7, 13)
BG_PANEL = RGBColor(11, 16, 32)
CARD = RGBColor(17, 24, 43)
CYAN = RGBColor(0, 240, 255)
PURPLE = RGBColor(138, 46, 255)
PINK = RGBColor(255, 0, 200)
WHITE = RGBColor(240, 245, 255)
MUTED = RGBColor(176, 188, 212)


def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return color[0], color[1], color[2]


def make_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    font = ImageFont.load_default()

    def new_canvas() -> Image.Image:
        return Image.new("RGB", (1280, 720), rgb_tuple(BG_PANEL))

    # Architecture image
    arch = new_canvas()
    d = ImageDraw.Draw(arch)
    d.text((30, 20), "CADGuard AI - System Pipeline", fill=rgb_tuple(CYAN), font=font)
    labels = ["User", "Upload", "Validate", "Simulate", "AI", "Output"]
    x = 40
    for idx, label in enumerate(labels):
        d.rounded_rectangle((x, 290, x + 170, 410), radius=18, outline=rgb_tuple(CYAN), width=3, fill=rgb_tuple(CARD))
        d.text((x + 58, 345), label, fill=rgb_tuple(WHITE), font=font)
        if idx < len(labels) - 1:
            d.line((x + 170, 350, x + 210, 350), fill=rgb_tuple(PURPLE), width=5)
            d.polygon([(x + 210, 350), (x + 195, 342), (x + 195, 358)], fill=rgb_tuple(PURPLE))
        x += 210
    arch_path = ASSET_DIR / "architecture.png"
    arch.save(arch_path)

    # Cost curve image
    curve = new_canvas()
    d = ImageDraw.Draw(curve)
    d.text((30, 20), "Late Error Detection -> Exponential Cost", fill=rgb_tuple(PINK), font=font)
    d.rectangle((80, 80, 1200, 620), outline=rgb_tuple(CYAN), width=2)
    points = [(120, 580), (240, 560), (360, 540), (500, 500), (650, 430), (820, 320), (980, 170), (1140, 95)]
    d.line(points, fill=rgb_tuple(PINK), width=7)
    for p in points:
        d.ellipse((p[0] - 8, p[1] - 8, p[0] + 8, p[1] + 8), fill=rgb_tuple(CYAN))
    d.text((1030, 130), "High Cost", fill=rgb_tuple(WHITE), font=font)
    curve_path = ASSET_DIR / "cost_curve.png"
    curve.save(curve_path)

    # Heatmap image
    heat = new_canvas()
    d = ImageDraw.Draw(heat)
    d.text((30, 20), "Stress Heatmap Preview", fill=rgb_tuple(CYAN), font=font)
    for i in range(0, 16):
        c = (min(255, 10 + i * 15), max(0, 210 - i * 10), min(255, 120 + i * 5))
        d.rectangle((130 + i * 62, 220, 180 + i * 62, 500), fill=c)
    d.rectangle((100, 190, 1180, 530), outline=rgb_tuple(WHITE), width=3)
    heat_path = ASSET_DIR / "heatmap.png"
    heat.save(heat_path)

    # Dashboard mock
    dash = new_canvas()
    d = ImageDraw.Draw(dash)
    d.text((30, 20), "CADGuard AI Dashboard Mock", fill=rgb_tuple(CYAN), font=font)
    d.rounded_rectangle((50, 80, 1230, 660), radius=20, outline=rgb_tuple(CYAN), width=3)
    # cards
    card_positions = [(90, 130), (380, 130), (670, 130), (960, 130)]
    for x, y in card_positions:
        d.rounded_rectangle((x, y, x + 220, y + 120), radius=16, outline=rgb_tuple(PURPLE), width=3, fill=rgb_tuple(CARD))
    d.rounded_rectangle((90, 290, 760, 620), radius=16, outline=rgb_tuple(CYAN), width=3, fill=rgb_tuple(CARD))
    d.rounded_rectangle((790, 290, 1190, 620), radius=16, outline=rgb_tuple(PINK), width=3, fill=rgb_tuple(CARD))
    dash_path = ASSET_DIR / "dashboard.png"
    dash.save(dash_path)

    # Chat mock
    chat = new_canvas()
    d = ImageDraw.Draw(chat)
    d.text((30, 20), "AI Copilot + Chatbot", fill=rgb_tuple(CYAN), font=font)
    d.rounded_rectangle((110, 100, 1170, 620), radius=18, outline=rgb_tuple(CYAN), width=3, fill=rgb_tuple(CARD))
    d.rounded_rectangle((150, 160, 620, 240), radius=12, outline=rgb_tuple(PURPLE), width=2, fill=rgb_tuple(BG_PANEL))
    d.rounded_rectangle((650, 260, 1120, 350), radius=12, outline=rgb_tuple(CYAN), width=2, fill=rgb_tuple(BG_PANEL))
    d.rounded_rectangle((150, 380, 700, 470), radius=12, outline=rgb_tuple(PINK), width=2, fill=rgb_tuple(BG_PANEL))
    d.text((180, 190), "Explain issue in this bracket", fill=rgb_tuple(WHITE), font=font)
    d.text((680, 300), "AI: Increase thickness to >=2.0mm", fill=rgb_tuple(WHITE), font=font)
    d.text((180, 415), "Auto-fix plan generated", fill=rgb_tuple(WHITE), font=font)
    chat_path = ASSET_DIR / "chat.png"
    chat.save(chat_path)

    # Roadmap image
    road = new_canvas()
    d = ImageDraw.Draw(road)
    d.text((30, 20), "Future Roadmap", fill=rgb_tuple(CYAN), font=font)
    d.line((120, 360, 1160, 360), fill=rgb_tuple(PURPLE), width=8)
    milestones = [
        (180, "Now", "Validation + AI"),
        (450, "Next", "FEM Module"),
        (730, "Later", "CAD Plugins"),
        (1010, "Scale", "Industry Suite"),
    ]
    for x, top, bottom in milestones:
        d.ellipse((x - 22, 338, x + 22, 382), fill=rgb_tuple(CYAN))
        d.text((x - 30, 300), top, fill=rgb_tuple(WHITE), font=font)
        d.text((x - 55, 395), bottom, fill=rgb_tuple(MUTED), font=font)
    road_path = ASSET_DIR / "roadmap.png"
    road.save(road_path)

    return {
        "architecture": arch_path,
        "cost_curve": curve_path,
        "heatmap": heat_path,
        "dashboard": dash_path,
        "chat": chat_path,
        "roadmap": road_path,
    }


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

    top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = CYAN
    top_line.line.fill.background()

    strip = slide.shapes.add_shape(1, Inches(12.95), Inches(0), Inches(0.38), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = PURPLE
    strip.line.fill.background()


def add_title_block(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.58), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.75), Inches(11.0), Inches(0.6))
        sub_tf = sub_box.text_frame
        sub_tf.clear()
        p2 = sub_tf.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = CYAN


def add_bullet_card(slide, title, bullets, x=0.7, y=2.1, w=7.4, h=4.8):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = CYAN
    card.line.width = Pt(1.5)

    tf = card.text_frame
    tf.clear()
    hdr = tf.paragraphs[0]
    hdr.text = title
    hdr.font.size = Pt(18)
    hdr.font.bold = True
    hdr.font.color.rgb = CYAN
    hdr.space_after = Pt(8)

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(21)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)


def add_kpi_row(slide, labels, x=0.75, y=5.95, w=11.8, h=0.95):
    segment = w / len(labels)
    for i, (label, value, color) in enumerate(labels):
        box = slide.shapes.add_shape(1, Inches(x + i * segment), Inches(y), Inches(segment - 0.12), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_PANEL
        box.line.color.rgb = color
        box.line.width = Pt(1.2)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER


def add_image_panel(slide, img_path: Path, caption: str, x=8.35, y=2.1, w=4.25, h=3.8):
    panel = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = BG_PANEL
    panel.line.color.rgb = CYAN
    panel.line.width = Pt(1.7)

    slide.shapes.add_picture(str(img_path), Inches(x + 0.12), Inches(y + 0.12), Inches(w - 0.24), Inches(h - 0.8))

    cap = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + h - 0.6), Inches(w - 0.3), Inches(0.45))
    tf = cap.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(12)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER


def add_highlight_line(slide, text):
    hl = slide.shapes.add_shape(1, Inches(0.74), Inches(6.85), Inches(12.1), Inches(0.42))
    hl.fill.solid()
    hl.fill.fore_color.rgb = BG_PANEL
    hl.line.color.rgb = PINK
    hl.line.width = Pt(1.2)
    tf = hl.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def add_section_footer(slide, text="CADGuard AI | Intelligent Design Validation Platform"):
    footer = slide.shapes.add_textbox(Inches(0.7), Inches(7.22), Inches(12.0), Inches(0.2))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


assets = make_assets()
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# 1 Title
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "CADGuard AI - Intelligent Design Validation Platform", "AI-powered CAD validation, simulation, and optimization")
add_bullet_card(s, "Presenter", ["Name: Your Name", "College: Your College Name", "Use case: Hackathon / Technical Review / Academic Evaluation"], x=0.8, y=2.35, w=6.3, h=2.7)
add_image_panel(s, assets["dashboard"], "Platform interface preview", x=7.35, y=2.15, w=5.3, h=4.3)
add_kpi_row(s, [("Accuracy-driven", "AI + Rules", CYAN), ("Fast", "Real-time", PURPLE), ("Actionable", "Fix Insights", PINK)])
add_notes(s, "Open with one sentence: CADGuard AI helps engineering teams catch and fix CAD issues early using AI and simulation.")
add_section_footer(s)

# 2 Problem
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "Problem")
add_bullet_card(s, "Current Industry Pain", ["Manual CAD validation is slow and inconsistent", "Expert-dependent review creates bottlenecks", "Late issue discovery causes redesign loops", "Prototype-stage defects inflate costs and delays"], x=0.75, y=2.0, w=7.4, h=4.6)
add_image_panel(s, assets["cost_curve"], "Cost grows rapidly with late error detection", x=8.3, y=2.0, w=4.4, h=4.6)
add_highlight_line(s, "Late design errors significantly increase product cost and time.")
add_notes(s, "Connect this to real project impact: rework, missed deadlines, and manufacturing waste.")
add_section_footer(s)

# 3 Solution
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "Our Solution")
add_bullet_card(s, "CADGuard AI Approach", ["AI-powered CAD validation engine", "Real-time geometry and rule analysis", "Stress simulation for risk preview", "AI Copilot recommendations for improvement", "Integrated reports and decision support"], x=0.75, y=2.0, w=7.5, h=4.6)
add_image_panel(s, assets["chat"], "AI assistant + recommendation workflow", x=8.35, y=2.0, w=4.3, h=4.6)
add_highlight_line(s, "An intelligent system that detects and improves designs early.")
add_notes(s, "Explain the value clearly: CADGuard AI is not only detection, it is guidance and optimization.")
add_section_footer(s)

# 4 Architecture
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "System Overview")
add_bullet_card(s, "Technology Stack", ["Frontend: React web platform", "Backend: FastAPI services", "Data: Supabase (DB + realtime)", "AI Layer: Groq API", "CAD Processing: trimesh geometry analysis"], x=0.75, y=2.0, w=7.4, h=4.6)
add_image_panel(s, assets["architecture"], "User -> Upload -> Validate -> Simulate -> AI -> Output", x=8.25, y=2.0, w=4.55, h=4.6)
add_notes(s, "Walk left to right across the pipeline and mention where each component contributes.")
add_section_footer(s)

# 5 Features
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "Key Features")
add_bullet_card(s, "Platform Capabilities", ["CAD Validation Engine", "Stress Simulation", "AI Design Copilot", "Chatbot Assistant", "3D Visualization", "Comparison Tool", "Report Generation (JSON/PDF)", "Live Notifications + Dashboard"], x=0.75, y=2.0, w=7.5, h=4.8)
add_image_panel(s, assets["dashboard"], "Interactive command-center UI", x=8.35, y=2.0, w=4.3, h=4.8)
add_notes(s, "Keep this brisk: each feature connects to a single practical benefit for the engineer.")
add_section_footer(s)

# 6 Demo
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "Demo Walkthrough")
add_bullet_card(s, "End-to-End Flow", ["1) Upload CAD model", "2) Run validation", "3) Review issues and severity", "4) Run simulation", "5) Get AI suggestions", "6) Compare two designs", "7) Export structured reports"], x=0.75, y=2.0, w=7.4, h=4.8)
add_image_panel(s, assets["dashboard"], "Live flow through product screens", x=8.3, y=2.0, w=4.35, h=4.8)
add_notes(s, "Narrate exactly what judges will see in the demo and emphasize speed of decision-making.")
add_section_footer(s)

# 7 Validation simulation
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "Technical Capabilities")
add_bullet_card(s, "Engineering Core", ["Rule-based deterministic validation", "Geometry quality checks", "Approximate stress simulation", "Weak-region identification", "Risk-level prediction"], x=0.75, y=2.0, w=7.35, h=4.8)
add_image_panel(s, assets["heatmap"], "Stress distribution heatmap preview", x=8.25, y=2.0, w=4.45, h=4.8)
add_highlight_line(s, "Uses deterministic rules and physics-based approximation.")
add_notes(s, "Be transparent and credible: this is a fast approximation engine designed for early-stage decision support.")
add_section_footer(s)

# 8 AI
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "AI Features")
add_bullet_card(s, "Intelligence Layer", ["Groq-powered AI Copilot", "Context-aware chatbot", "Error-fixing assistant", "Human-readable explanations", "Actionable design changes"], x=0.75, y=2.0, w=7.4, h=4.7)
add_image_panel(s, assets["chat"], "AI conversation + fix recommendations", x=8.3, y=2.0, w=4.35, h=4.7)
add_highlight_line(s, "AI converts complex engineering data into actionable insights.")
add_notes(s, "Position AI as an engineering assistant: it explains and guides, while deterministic checks keep baseline reliability.")
add_section_footer(s)

# 9 UX
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "User Experience")
add_bullet_card(s, "Design Philosophy", ["Neon futuristic UI", "Interactive analytics dashboard", "Realtime updates and alerts", "3D CAD + stress visualization", "Fast navigation for engineering workflows"], x=0.75, y=2.0, w=7.35, h=4.7)

# three screenshot cards
for i, label in enumerate(["Dashboard", "Validation", "Compare"]):
    x = 8.2 + i * 1.45
    card = s.shapes.add_shape(1, Inches(x), Inches(2.4), Inches(1.35), Inches(2.45))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_PANEL
    card.line.color.rgb = [CYAN, PURPLE, PINK][i]
    card.line.width = Pt(1.2)
    s.shapes.add_picture(str(assets["dashboard"]), Inches(x + 0.05), Inches(2.45), Inches(1.25), Inches(1.9))
    t = s.shapes.add_textbox(Inches(x), Inches(4.95), Inches(1.35), Inches(0.4))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER

add_notes(s, "Explain that UX is optimized for clarity and confidence, making advanced analysis usable in a fast-paced review cycle.")
add_section_footer(s)

# 10 uniqueness
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "What Makes It Different?")
add_bullet_card(s, "Differentiators", ["Combines CAD + AI + Simulation in one platform", "Beyond issue detection: provides fix suggestions", "Realtime workflow from upload to report", "Built for iterative engineering improvements"], x=0.75, y=2.0, w=7.35, h=4.8)
add_image_panel(s, assets["architecture"], "Unified workflow instead of isolated tools", x=8.25, y=2.0, w=4.45, h=4.8)
add_notes(s, "Differentiate from single-point tools: CADGuard AI closes the loop from detection to recommendation and action.")
add_section_footer(s)

# 11 future
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "Future Enhancements")
add_bullet_card(s, "Roadmap", ["Real FEM simulation integration", "Native CAD plugins", "ML-based predictive defect models", "Enterprise manufacturing integrations", "Expanded materials and load libraries"], x=0.75, y=2.0, w=7.35, h=4.8)
add_image_panel(s, assets["roadmap"], "Scalable roadmap to production usage", x=8.25, y=2.0, w=4.45, h=4.8)
add_notes(s, "Show ambition and realism: immediate technical depth plus long-term productization path.")
add_section_footer(s)

# 12 conclusion
s = prs.slides.add_slide(blank)
set_background(s)
add_title_block(s, "Conclusion")
add_bullet_card(s, "Closing", ["Early detection reduces redesign cost", "AI guidance speeds engineering decisions", "Simulation adds risk visibility", "CADGuard AI moves validation from manual to intelligent"], x=0.8, y=2.1, w=7.15, h=4.2)
add_image_panel(s, assets["dashboard"], "Thank you - Questions?", x=8.2, y=2.1, w=4.5, h=4.2)
closing = s.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.2), Inches(0.45))
ctf = closing.text_frame
ctf.clear()
p = ctf.paragraphs[0]
p.text = "CADGuard AI transforms design validation from manual to intelligent."
p.font.size = Pt(19)
p.font.bold = True
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER
add_notes(s, "Close confidently: summarize impact and invite questions.")
add_section_footer(s)

prs.save(str(OUTPUT_PATH))
print(f"Presentation created: {OUTPUT_PATH}")
